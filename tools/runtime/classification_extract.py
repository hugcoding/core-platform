#!/usr/bin/env python3
"""Locally extract and classify modern golden documents without retaining raw text."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import signal
import sys
import time
from collections import Counter
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from zoneinfo import ZoneInfo

from core.exports.csv_format import dict_reader, write_dict_rows


CATEGORY_TERMS = {
    "gevoelig/financiën": (
        "belasting", "aanslag", "bankrekening", "rekeningnummer", "hypotheek",
        "verzekering", "jaaropgave", "pensioen", "financieel",
    ),
    "gevoelig/gezondheid": (
        "huisarts", "ziekenhuis", "diagnose", "medisch", "medicatie",
        "gezondheid", "behandeling", "zorgverlener", "wmo",
    ),
    "gevoelig/identiteit": (
        "paspoort", "rijbewijs", "identiteitskaart", "burgerservicenummer", "bsn",
    ),
    "gevoelig/werk_en_inkomen": (
        "arbeidsovereenkomst", "werkgever", "salaris", "loonstrook",
        "sollicitatie", "curriculum vitae", "functioneringsgesprek",
    ),
    "documenten/studie": (
        "opleiding", "studie", "module", "les", "opdracht", "examen",
        "scriptie", "hoofdstuk", "student", "docent", "ncoi",
    ),
    "documenten/werk": (
        "projectplan", "vergadering", "notulen", "werkoverleg", "organisatie",
        "management", "proces", "klant", "offerte",
    ),
    "documenten/wonen": (
        "woning", "verbouwing", "energieleverancier", "gevel", "keuken",
        "badkamer", "huur", "koopakte", "kadaster",
    ),
    "documenten/administratie": (
        "factuur", "aankoop", "garantie", "abonnement", "overeenkomst",
        "bevestiging", "betaling", "bon", "nota",
    ),
    "documenten/persoonlijk": (
        "familie", "verjaardag", "vakantie", "persoonlijk", "uitnodiging",
        "wensenlijst", "reis", "ticket",
    ),
    "projecten": (
        "source code", "database", "software", "server", "netwerk",
        "openvpn", "uml", "xml", "api", "configuratie",
    ),
}

EXTRACTION_VERSION = "temporal-v2"
DEFAULT_TIMEOUT_SECONDS = 120
DEFAULT_CHECKPOINT_EVERY = 10


class DocumentExtractionTimeout(TimeoutError):
    """Raised when one document exceeds its configured processing deadline."""


@contextmanager
def extraction_deadline(seconds: int):
    """Enforce a per-document deadline in the Linux classifier container."""
    if seconds <= 0 or not hasattr(signal, "SIGALRM"):
        yield
        return

    def handle_timeout(_signum, _frame):
        raise DocumentExtractionTimeout(
            f"document processing exceeded {seconds} seconds"
        )

    previous_handler = signal.signal(signal.SIGALRM, handle_timeout)
    signal.setitimer(signal.ITIMER_REAL, seconds)
    try:
        yield
    finally:
        signal.setitimer(signal.ITIMER_REAL, 0)
        signal.signal(signal.SIGALRM, previous_handler)


def normalize(text: str) -> str:
    return " ".join(text.split()).casefold()


def iso_value(value) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


def parse_loose_pdf_date(value: object) -> str:
    """Parse common non-standard PDF date values without failing text extraction."""
    raw = str(value or "").strip()
    if not raw:
        return ""
    candidate = raw[2:] if raw.startswith("D:") else raw
    for pattern in ("%m/%d/%Y %H:%M:%S", "%m/%d/%Y %H:%M:%S%z"):
        try:
            return datetime.strptime(candidate, pattern).isoformat()
        except ValueError:
            pass
    match = re.fullmatch(
        r"(?P<year>\d{4})(?P<month>\d{2})(?P<day>\d{2})"
        r"(?P<hour>\d{2})(?P<minute>\d{2})(?P<second>\d{2})"
        r"(?:(?P<z>Z)|(?P<sign>[+-])(?P<offhour>\d{2})'?"
        r"(?P<offminute>\d{2})'?)?",
        candidate,
    )
    if not match:
        raise ValueError(f"unsupported PDF date: {raw}")
    values = {key: int(match.group(key)) for key in (
        "year", "month", "day", "hour", "minute", "second"
    )}
    tz = None
    if match.group("z"):
        tz = timezone.utc
    elif match.group("sign"):
        offset_minutes = int(match.group("offhour")) * 60 + int(match.group("offminute"))
        if match.group("sign") == "-":
            offset_minutes *= -1
        from datetime import timedelta
        tz = timezone(timedelta(minutes=offset_minutes))
    return datetime(**values, tzinfo=tz).isoformat()


def pdf_metadata_evidence(metadata) -> dict[str, object]:
    """Read PDF metadata defensively and retain warnings separately from content."""
    evidence: dict[str, object] = {}
    warnings: list[str] = []
    for output_name, attribute, raw_key in (
        ("embedded_created_at", "creation_date", "/CreationDate"),
        ("embedded_modified_at", "modification_date", "/ModDate"),
    ):
        try:
            evidence[output_name] = iso_value(getattr(metadata, attribute, None))
        except (TypeError, ValueError) as exc:
            try:
                evidence[output_name] = parse_loose_pdf_date(metadata.get(raw_key))
            except (TypeError, ValueError):
                evidence[output_name] = ""
                warnings.append(f"{output_name}: {type(exc).__name__}: {exc}")
    try:
        evidence["embedded_author"] = iso_value(getattr(metadata, "author", None))
    except (TypeError, ValueError) as exc:
        evidence["embedded_author"] = ""
        warnings.append(f"embedded_author: {type(exc).__name__}: {exc}")
    evidence["metadata_warnings"] = warnings
    return evidence


def extract_text(path: Path, route: str) -> tuple[str, int | None, dict[str, object]]:
    if route == "plain-text":
        try:
            return path.read_text(encoding="utf-8"), None, {}
        except UnicodeDecodeError:
            return path.read_text(encoding="cp1252"), None, {}
    if route == "pypdf":
        from pypdf import PdfReader

        reader = PdfReader(path)
        if reader.is_encrypted:
            if reader.decrypt("") == 0:
                raise PermissionError("PDF requires a password")
        metadata = reader.metadata or {}
        page_text: list[str] = []
        content_warnings: list[str] = []
        for page_number, page in enumerate(reader.pages, start=1):
            try:
                page_text.append(page.extract_text() or "")
            except Exception as exc:
                page_text.append("")
                content_warnings.append(
                    f"page_{page_number}: {type(exc).__name__}: {exc}"
                )
        evidence = pdf_metadata_evidence(metadata)
        evidence["content_warnings"] = content_warnings
        return "\n".join(page_text), len(reader.pages), evidence
    if route == "python-docx":
        from docx import Document

        document = Document(path)
        blocks = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                blocks.append(" ".join(cell.text for cell in row.cells))
        properties = document.core_properties
        return "\n".join(blocks), None, {
            "embedded_created_at": iso_value(properties.created),
            "embedded_modified_at": iso_value(properties.modified),
            "embedded_author": iso_value(properties.last_modified_by or properties.author),
            "embedded_revision": iso_value(properties.revision),
        }
    if route == "openpyxl":
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)
        values = []
        for sheet in workbook.worksheets:
            values.append(sheet.title)
            for row in sheet.iter_rows(values_only=True):
                values.extend(str(value) for value in row if value is not None)
        properties = workbook.properties
        metadata = {
            "embedded_created_at": iso_value(properties.created),
            "embedded_modified_at": iso_value(properties.modified),
            "embedded_author": iso_value(properties.lastModifiedBy or properties.creator),
        }
        workbook.close()
        return "\n".join(values), None, metadata
    if route == "python-pptx":
        from pptx import Presentation

        presentation = Presentation(path)
        values = []
        for slide in presentation.slides:
            values.extend(
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
        properties = presentation.core_properties
        return "\n".join(values), len(presentation.slides), {
            "embedded_created_at": iso_value(properties.created),
            "embedded_modified_at": iso_value(properties.modified),
            "embedded_author": iso_value(properties.last_modified_by or properties.author),
            "embedded_revision": iso_value(properties.revision),
        }
    if route == "rtf":
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(path.read_text(encoding="cp1252", errors="replace")), None, {}
    if route == "odf":
        from odf import teletype
        from odf.opendocument import load

        document = load(str(path))
        return teletype.extractText(document), None, {}
    raise ValueError(f"unsupported extraction route: {route}")


def term_hits(text: str, terms: tuple[str, ...]) -> list[str]:
    return [term for term in terms if re.search(rf"\b{re.escape(term)}\b", text)]


DATE_PATTERNS = (
    r"\b(?:19|20)\d{2}[-/.](?:0?[1-9]|1[0-2])[-/.](?:0?[1-9]|[12]\d|3[01])\b",
    r"\b(?:0?[1-9]|[12]\d|3[01])[-/.](?:0?[1-9]|1[0-2])[-/.](?:19|20)\d{2}\b",
)


def date_candidates(text: str, limit: int = 50) -> list[str]:
    values = []
    for pattern in DATE_PATTERNS:
        values.extend(re.findall(pattern, text))
    return sorted(set(values))[:limit]


LOCAL_TIMEZONE = ZoneInfo("Europe/Amsterdam")


def parse_embedded_datetime(value: str) -> datetime | None:
    """Parse embedded document time and localize offset-less wall-clock values."""
    if not value:
        return None
    parsed = datetime.fromisoformat(value.replace("Z", "+00:00"))
    if parsed.tzinfo is None:
        parsed = parsed.replace(tzinfo=LOCAL_TIMEZONE)
    return parsed


def temporal_inconsistencies(created: str, modified: str) -> list[str]:
    issues = []
    try:
        created_value = parse_embedded_datetime(created)
        modified_value = parse_embedded_datetime(modified)
        now = datetime.now(timezone.utc)
        for label, value in (("created", created_value), ("modified", modified_value)):
            if value:
                if value.astimezone(timezone.utc) > now:
                    issues.append(f"{label}_in_future")
        if created_value and modified_value:
            left = created_value.astimezone(timezone.utc)
            right = modified_value.astimezone(timezone.utc)
            if left > right:
                created_wall_time = created_value.replace(tzinfo=None)
                modified_wall_time = modified_value.replace(tzinfo=None)
                if (
                    created_value.utcoffset() != modified_value.utcoffset()
                    and created_wall_time <= modified_wall_time
                ):
                    issues.append("embedded_timezone_conflict")
                else:
                    issues.append("created_after_modified")
    except (TypeError, ValueError):
        issues.append("embedded_date_parse_error")
    return issues


def classify_content(text: str, filename: str, source_path: str) -> tuple[str, str, str]:
    normalized = normalize(text)
    name = normalize(filename)
    path = normalize(source_path)
    scores = {}
    evidence = {}
    for category, terms in CATEGORY_TERMS.items():
        content_hits = term_hits(normalized, terms)
        name_hits = term_hits(name, terms)
        path_hits = term_hits(path, terms)
        score = len(content_hits) * 5 + len(name_hits) * 2 + len(path_hits)
        scores[category] = score
        evidence[category] = {
            "inhoud": content_hits,
            "bestandsnaam": name_hits,
            "bronpad": path_hits,
        }
    ordered = sorted(scores, key=lambda category: (-scores[category], category))
    best = ordered[0]
    margin = scores[best] - scores[ordered[1]]
    if scores[best] == 0:
        return "documenten/uitzoeken", "low", "geen inhoudssignalen"
    confidence = "high" if scores[best] >= 10 and margin >= 3 else "medium" if margin > 0 else "low"
    reasons = evidence[best]
    return best, confidence, (
        f"score={scores[best]}; margin={margin}; "
        f"inhoud={','.join(reasons['inhoud']) or '-'}; "
        f"bestandsnaam={','.join(reasons['bestandsnaam']) or '-'}; "
        f"bronpad={','.join(reasons['bronpad']) or '-'}"
    )


def process_row(
    row: dict[str, str], timeout_seconds: int = DEFAULT_TIMEOUT_SECONDS
) -> dict[str, str]:
    started = time.monotonic()
    result = {
        **row,
        "extraction_result": "",
        "extraction_error": "",
        "extraction_warnings": "[]",
        "metadata_extraction_status": "not_attempted",
        "extraction_elapsed_seconds": "0.000",
        "characters": "",
        "words": "",
        "pages_or_slides": "",
        "filesystem_mtime": "",
        "embedded_created_at": "",
        "embedded_modified_at": "",
        "embedded_author_present": "",
        "embedded_revision": "",
        "filename_date_candidates": "[]",
        "content_date_candidates": "[]",
        "temporal_inconsistencies": "[]",
        "temporal_assessment_status": "evidence_only",
        "temporal_extraction_version": EXTRACTION_VERSION,
    }
    if row["extraction_status"] != "ready_for_local_extraction":
        result["extraction_result"] = "skipped"
        return result
    if row.get("size_bytes") == "0":
        category, confidence, reasons = classify_content(
            "", row["filename"], row["golden_path"]
        )
        result.update(
            {
                "extraction_result": "empty_file",
                "content_category": category,
                "category_confidence": confidence,
                "category_reasons": reasons,
                "characters": "0",
                "words": "0",
            }
        )
        return result
    try:
        path = Path(row["golden_path"])
        with extraction_deadline(timeout_seconds):
            text, pages, embedded = extract_text(path, row["extraction_route"])
        normalized = normalize(text)
        category, confidence, reasons = classify_content(
            normalized, row["filename"], row["golden_path"]
        )
        needs_ocr = row["extraction_route"] == "pypdf" and not normalized and bool(pages)
        content_warnings = embedded.get("content_warnings", [])
        metadata_warnings = embedded.get("metadata_warnings", [])
        extraction_result = (
            "partial_extraction"
            if content_warnings and normalized
            else "needs_ocr" if needs_ocr else "extracted"
        )
        result.update(
            {
                "extraction_result": extraction_result,
                "extraction_error": "",
                "extraction_warnings": json.dumps(
                    [*content_warnings, *metadata_warnings], ensure_ascii=False
                ),
                "metadata_extraction_status": (
                    "metadata_parse_warning"
                    if metadata_warnings else "extracted"
                ),
                "characters": str(len(normalized)),
                "words": str(len(normalized.split())) if normalized else "0",
                "pages_or_slides": str(pages or ""),
                "filesystem_mtime": datetime.fromtimestamp(
                    path.stat().st_mtime, tz=timezone.utc
                ).isoformat(),
                "embedded_created_at": str(embedded.get("embedded_created_at", "")),
                "embedded_modified_at": str(embedded.get("embedded_modified_at", "")),
                "embedded_author_present": str(bool(embedded.get("embedded_author"))).lower(),
                "embedded_revision": embedded.get("embedded_revision", ""),
                "filename_date_candidates": json.dumps(
                    date_candidates(row["filename"]), ensure_ascii=False
                ),
                "content_date_candidates": json.dumps(
                    date_candidates(text), ensure_ascii=False
                ),
                "temporal_inconsistencies": json.dumps(
                    temporal_inconsistencies(
                        str(embedded.get("embedded_created_at", "")),
                        str(embedded.get("embedded_modified_at", "")),
                    ),
                    ensure_ascii=False,
                ),
                "content_category": category,
                "category_confidence": confidence,
                "category_reasons": reasons,
                "proposed_target_path": "",
            }
        )
    except DocumentExtractionTimeout as exc:
        result.update(
            {
                "extraction_result": "extraction_timeout",
                "extraction_error": str(exc),
                "characters": "0",
                "words": "0",
                "pages_or_slides": "",
            }
        )
    except Exception as exc:
        result.update(
            {
                "extraction_result": (
                    "password_required" if isinstance(exc, PermissionError) else "error"
                ),
                "extraction_error": f"{type(exc).__name__}: {exc}",
                "characters": "0",
                "words": "0",
                "pages_or_slides": "",
            }
        )
    finally:
        result["extraction_elapsed_seconds"] = f"{time.monotonic() - started:.3f}"
    return result


def resolve_manifest(root: Path, value: str) -> Path:
    export_dir = root / "project" / "exports" / "migration-inventory"
    if value == "latest":
        manifests = sorted(export_dir.glob("classification-inventory-*.csv"))
        if not manifests:
            raise FileNotFoundError("No classification inventory found.")
        return manifests[-1]
    path = Path(value)
    return path if path.is_absolute() else root / path


def row_identity(row: dict[str, str]) -> str:
    return row.get("content_group_id") or row.get("golden_file_id", "")


def checkpoint_path(export_dir: Path, manifest_path: Path) -> Path:
    digest = hashlib.sha256(manifest_path.read_bytes()).hexdigest()[:12]
    return export_dir / f".classification-extract-{digest}-{EXTRACTION_VERSION}.csv"


def read_checkpoint(path: Path) -> list[dict[str, str]]:
    if not path.exists():
        return []
    with path.open(newline="", encoding="utf-8-sig") as handle:
        return list(dict_reader(handle))


def write_checkpoint(path: Path, results: list[dict[str, str]]) -> None:
    if results:
        temporary = path.with_suffix(path.suffix + ".tmp")
        write_dict_rows(temporary, results, list(results[0]))
        temporary.replace(path)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract and classify golden documents locally.")
    parser.add_argument("--manifest", default="latest")
    parser.add_argument("--dry-run", action="store_true", required=True)
    parser.add_argument("--timeout-seconds", type=int, default=DEFAULT_TIMEOUT_SECONDS)
    parser.add_argument("--checkpoint-every", type=int, default=DEFAULT_CHECKPOINT_EVERY)
    parser.add_argument("--progress-every", type=int, default=25)
    parser.add_argument(
        "--resume", action=argparse.BooleanOptionalAction, default=True,
        help="resume the matching manifest/extractor checkpoint (default: enabled)",
    )
    args = parser.parse_args(argv)
    if args.timeout_seconds < 0 or args.checkpoint_every < 1 or args.progress_every < 1:
        parser.error("timeout must be >= 0; checkpoint/progress intervals must be >= 1")
    root = Path("/app")
    try:
        manifest_path = resolve_manifest(root, args.manifest)
        with manifest_path.open(newline="", encoding="utf-8-sig") as handle:
            rows = list(dict_reader(handle))
    except (FileNotFoundError, OSError) as exc:
        print(f"Classification extraction failed: {exc}")
        return 1

    export_dir = root / "project" / "exports" / "migration-inventory"
    checkpoint = checkpoint_path(export_dir, manifest_path)
    results = read_checkpoint(checkpoint) if args.resume else []
    completed = {row_identity(row) for row in results}
    resumed_count = len(results)
    total = len(rows)
    batch_started = time.monotonic()
    if resumed_count:
        print(
            f"Resuming {resumed_count}/{total} records from checkpoint "
            f"{checkpoint.name}",
            flush=True,
        )

    newly_processed = 0
    for position, row in enumerate(rows, start=1):
        if row_identity(row) in completed:
            continue
        item_started = time.monotonic()
        result = process_row(row, timeout_seconds=args.timeout_seconds)
        results.append(result)
        completed.add(row_identity(row))
        newly_processed += 1
        elapsed = time.monotonic() - item_started
        print(
            f"[{position}/{total}] file_id={row.get('golden_file_id', '-')} "
            f"route={row.get('extraction_route', '-')} "
            f"status={result['extraction_result']} elapsed={elapsed:.1f}s",
            flush=True,
        )
        if newly_processed % args.checkpoint_every == 0:
            write_checkpoint(checkpoint, results)
        if newly_processed % args.progress_every == 0:
            interim = Counter(item.get("extraction_result", "") for item in results)
            print(
                f"Progress: processed={len(results)} remaining={total - len(results)} "
                f"errors={interim['error']} timeouts={interim['extraction_timeout']} "
                f"elapsed={time.monotonic() - batch_started:.1f}s",
                file=sys.stderr,
                flush=True,
            )
    write_checkpoint(checkpoint, results)

    timestamp = datetime.now().astimezone().strftime("%Y%m%d-%H%M%S")
    output_path = export_dir / f"classification-results-{timestamp}.csv"
    report_path = export_dir / f"classification-results-{timestamp}.md"
    write_dict_rows(output_path, results, list(results[0]))
    outcomes = Counter(row.get("extraction_result", "") for row in results)
    categories = Counter(row.get("content_category", "") for row in results)
    embedded_created = sum(bool(row.get("embedded_created_at")) for row in results)
    embedded_modified = sum(bool(row.get("embedded_modified_at")) for row in results)
    content_dates = sum(row.get("content_date_candidates", "[]") != "[]" for row in results)
    temporal_issues = sum(row.get("temporal_inconsistencies", "[]") != "[]" for row in results)
    metadata_warnings = sum(
        row.get("metadata_extraction_status") == "metadata_parse_warning"
        for row in results
    )
    total_elapsed = time.monotonic() - batch_started
    slowest = sorted(
        results,
        key=lambda row: float(row.get("extraction_elapsed_seconds") or 0),
        reverse=True,
    )[:10]
    report = [
        "# SCRUM-61 lokale inhoudsclassificatie", "",
        f"- Gegenereerd: `{datetime.now().astimezone().isoformat()}`",
        "- Modus: **lokaal, alleen-lezen, zonder embeddings**",
        f"- Extracted: **{outcomes['extracted']}**",
        f"- Gedeeltelijk geëxtraheerd: **{outcomes['partial_extraction']}**",
        f"- OCR vereist: **{outcomes['needs_ocr']}**",
        f"- Conversie/overgeslagen: **{outcomes['skipped']}**",
        f"- Lege bestanden: **{outcomes['empty_file']}**",
        f"- Wachtwoord vereist: **{outcomes['password_required']}**",
        f"- Timeouts: **{outcomes['extraction_timeout']}**",
        f"- Fouten: **{outcomes['error']}**",
        f"- Metadatawaarschuwingen: **{metadata_warnings}**",
        f"- Hervat uit checkpoint: **{resumed_count}**",
        f"- Timeout per document: **{args.timeout_seconds} seconden**",
        f"- Batchdoorlooptijd deze run: **{total_elapsed:.1f} seconden**",
        f"- Extractorversie: **{EXTRACTION_VERSION}**",
        f"- Ingebedde aanmaakdatum: **{embedded_created}**",
        f"- Ingebedde wijzigingsdatum: **{embedded_modified}**",
        f"- Inhoudelijke datumkandidaten: **{content_dates}**",
        f"- Tijdsinconsistenties: **{temporal_issues}**",
        "", "## Voorgestelde categorieën", "", "| Categorie | Bestanden |", "|---|---:|",
    ]
    report.extend(
        f"| `{category}` | {count} |"
        for category, count in sorted(categories.items())
        if category
    )
    report.extend([
        "", "## Langzaamste documenten", "",
        "| File ID | Route | Status | Seconden |", "|---:|---|---|---:|",
    ])
    report.extend(
        f"| `{row.get('golden_file_id', '-')}` | `{row.get('extraction_route', '-')}` "
        f"| `{row.get('extraction_result', '-')}` "
        f"| {float(row.get('extraction_elapsed_seconds') or 0):.3f} |"
        for row in slowest
    )
    report.extend([
        "", "Volledige documenttekst is niet opgeslagen.",
        "Er zijn geen bestanden, mappen of databaserecords gewijzigd.",
    ])
    report_path.write_text("\n".join(report) + "\n", encoding="utf-8")
    checkpoint.unlink(missing_ok=True)
    print("SCRUM-61 local classification extraction complete")
    print(f"Report: {report_path.relative_to(root)}")
    print(f"Results: {output_path.relative_to(root)}")
    return 0 if not outcomes["error"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
